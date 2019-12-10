import os
# pip install pillow
from PIL import Image, ImageTk
# pip install python-pptx
import pptx
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.dml.color import RGBColor


# Creates the window to be displayed
def createView(menuItems):
    prs = Presentation('resources/SpecialsSlide.pptx')
    slideShapes = prs.slides[0].shapes
    for i in range(5, 17):
        if i < 14:
            if (i - 5) % 3 == 0:
                slideShapes[i].text = menuItems[(i - 5)//3]['name']
                slideShapes[i].text_frame.paragraphs[0].font.bold = True
                slideShapes[i].text_frame.paragraphs[0].font.size = Pt(18)
            elif (i - 5) % 3 == 1:
                slideShapes[i].text = menuItems[(i - 5)//3]['description']
                slideShapes[i].text_frame.paragraphs[0].font.italic = True
                slideShapes[i].text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(
                    '737373')
                slideShapes[i].text_frame.paragraphs[0].font.size = Pt(16)
            else:
                slideShapes[i].text = menuItems[(i - 5)//3]['price']
                slideShapes[i].text_frame.paragraphs[0].font.bold = True
                slideShapes[i].text_frame.paragraphs[0].font.size = Pt(18)
        else:
            itemImage = Image.open('images/{}.gif'.format((i - 15) % 3 + 1))
            slideShapes[i].Image = itemImage
    prs.save('resources/TodaysSpecials.pptx')
    return
